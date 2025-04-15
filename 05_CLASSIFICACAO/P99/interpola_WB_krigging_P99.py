import numpy as np
import pandas as pd
import geopandas as gpd
import matplotlib.pyplot as plt
from pykrige.ok import OrdinaryKriging
from matplotlib.colors import ListedColormap, BoundaryNorm
from pptx import Presentation
from pptx.util import Cm
import time
import os
import xarray as xr

# Configurações iniciais
shapefile_path = '../../00_STUFF/BAHIA.shp'
gdf = gpd.read_file(shapefile_path)

# Definição pelo usuário
x_min, x_max, y_min, y_max, res = -47.9, -37, -19.9, -8, 0.1

# Definição pelos limites do shapefile
#x_min, y_min, x_max, y_max = gdf.total_bounds
print(x_min, x_max, y_min, y_max)
res = 0.1
gridx, gridy = np.arange(x_min, x_max, res), np.arange(y_min, y_max, res)

# Definições para Kriging com múltiplos variogramas
variogram_models = {
    "linear": {'nugget': 0.01, 'slope': 1.0},
    "power": {'nugget': 0.01, 'exponent': 1.0, 'scale': 1.0},
    "gaussian": {'nugget': 0.01, 'sill': 2.0, 'range': 2.0},
    "spherical": {'nugget': 0.01, 'sill': 1.0, 'range': 10.0},
    "exponential": {'nugget': 0.01, 'sill': 2.0, 'range': 10.0},
    "hole-effect": {'nugget': 0.01, 'sill': 2.0, 'range': 10.0},
}

# Definição dos intervalos e cores correspondentes
#intervalos = [10, 12, 14,16, 18, 20, 22, 24, 26, 28, 30, 32, 34, 36, 38, 40, 42, 44, 46, 48, 50]
intervalos = [40, 45, 50, 55 , 60 , 65, 70, 72, 74, 76, 78, 80, 82, 84, 86, 88, 90, 94, 96, 98, 100]

cores = [
    (1.0, 0.0, 0.0),  # Vermelho
    (1.0, 0.3, 0.0),  # Laranja avermelhado
    (1.0, 0.5, 0.0),  # Laranja
    (1.0, 0.7, 0.0),  # Laranja amarelado
    (1.0, 1.0, 0.0),  # Amarelo
    (0.8, 1.0, 0.2),  # Verde amarelado
    (0.6, 1.0, 0.4),  # Verde claro
    (0.0, 1.0, 0.0),  # Verde
    (0.0, 0.8, 0.4),  # Verde escuro
    #(0.3, 1.0, 0.6),  # Verde água
    'seagreen',
    (0.0, 0.6, 0.8),  # Azul esverdeado
    (0.0, 0.4, 1.0),  # Azul claro
    (0.0, 0.2, 1.0),  # Azul médio
    (0.0, 0.0, 1.0),  # Azul
    (0.3, 0.0, 1.0),  # Azul arroxeado
    (0.5, 0.0, 1.0),  # Roxo azulado
    (0.7, 0.0, 1.0),  # Roxo médio
    (0.8, 0.2, 1.0),  # Roxo claro
    (0.9, 0.4, 1.0),  # Lilás
    (1.0, 0.6, 1.0)   # Rosa arroxeado
]

# cores = ['darkred','brown','red', 'orangered', 'darkorange', 'orange', 'gold', 'yellow', 
         # 'lime', 'limegreen', 'green',  'teal', 'deepskyblue', 
         # 'blue', 'mediumblue', 'darkblue', 'blueviolet', 'mediumorchid','plum' ]
# Criando o mapa de cores e normalização
cmap = ListedColormap(cores)
norm = BoundaryNorm(intervalos, cmap.N)





# Função para interpolação Kriging com múltiplos variogramas
def interpola_kriging(olat, olon, valores, variogram_model, gridx, gridy, variogram_params):
    OK = OrdinaryKriging(
        olon,
        olat,
        valores,
        variogram_model=variogram_model,
        variogram_parameters=variogram_params,
        verbose=False,
        enable_plotting=False,
    )
    z, _ = OK.execute("grid", gridx, gridy)
    return z

# Função para salvar os dados interpolados em um arquivo NetCDF
def salvar_netcdf(filename, gridx, gridy, z):
    ds = xr.Dataset(
        {"Interpolacao": (("lat", "lon"), z)},
        coords={"lat": gridy, "lon": gridx},
    )
    ds.to_netcdf(filename)

# Função para plotar o mapa da interpolação
def plota_mapa(caso, gridx, gridy, z, variavel, intervalos, metodo, cmap):
    plt.figure(figsize=(10, 6))
    grid_x, grid_y = np.meshgrid(gridx, gridy)
    norm = BoundaryNorm(intervalos, cmap.N)
    plt.contourf(grid_x, grid_y, z, levels=intervalos, cmap=cmap, norm=norm)
    plt.xlabel('Longitude')
    plt.ylabel('Latitude')
    plt.xlim(x_min, x_max)
    plt.ylim(y_min, y_max)
    plt.title(f"{caso} - Interpolação {metodo}")
    gdf.plot(ax=plt.gca(), edgecolor='gray', facecolor='none', lw=1) 
    cbar = plt.colorbar(label="Média mm/ano")
    cbar.ax.yaxis.set_ticks(intervalos)  # Define os intervalos explicitamente
    cbar.ax.yaxis.set_ticklabels(intervalos)  # Define as labels explicitamente
    filename = f"{caso}_interpolacao_{metodo}.png"
    plt.savefig(filename, dpi=300, bbox_inches="tight")
    plt.close()
    return filename

# Função para gerar interpolação Kriging e criar apresentação PowerPoint
def gera_interpolacao_kriging(prefixo_imagem, arquivo_pptx):
    generated_images = []

    for model, params in variogram_models.items():
        z_kriging = interpola_kriging(olat, olon, valores, model, gridx, gridy, params)
        img_filename = plota_mapa(prefixo_imagem, gridx, gridy, z_kriging, "Média mm/ano", intervalos, f"Kriging_{model}", cmap)
        netcdf_filename = img_filename.replace(".png", ".nc")
        salvar_netcdf(netcdf_filename, gridx, gridy, z_kriging)
        generated_images.append(img_filename)

    prs = Presentation()
    prs.slide_width = Cm(29.7)
    prs.slide_height = Cm(21.0)
    slide_layout = prs.slide_layouts[5]
    escala = 0.65

    for i in range(0, len(generated_images), 6):
        slide = prs.slides.add_slide(slide_layout)
        for j in range(6):
            if i + j < len(generated_images):
                img_path = generated_images[i + j]
                left = Cm(0.5 + (j % 3) * 10.00)
                top = Cm(0.5 + (j // 3) * 10.00)
                slide.shapes.add_picture(img_path, left, top, width=Cm(14.03 * escala), height=Cm(12.65 * escala))

    prs.save(arquivo_pptx)
    print("Apresentação PowerPoint criada com sucesso!")

# Processamento dos diferentes períodos
datasets = [
    ("../../04_FILTRAGEM/selecao_BAHIA.xlsx", "bahia_P99_inventario.pptx", "bahia_inventario_"),
    ("../../04_FILTRAGEM/selecao_BAHIA_1981a2010.xlsx", "bahia_P99_1981a2010.pptx", "bahia_1981a2010_"),
    ("../../04_FILTRAGEM/selecao_BAHIA_1991a2020.xlsx", "bahia_P99_1991a2020.pptx", "bahia_1991a2020_"),
]

for file_path, arquivo_pptx, prefixo_imagem in datasets:
    df = pd.read_excel(file_path, sheet_name="ANUAL")
    olat, olon, valores = df["Latitude"], df["Longitude"], df["P99 NZ"]
    gera_interpolacao_kriging(prefixo_imagem, arquivo_pptx)
