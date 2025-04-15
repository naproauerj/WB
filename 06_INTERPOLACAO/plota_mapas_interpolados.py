import numpy as np
import geopandas as gpd
import matplotlib.pyplot as plt
import xarray as xr
from matplotlib.colors import ListedColormap, BoundaryNorm
from pptx import Presentation
from pptx.util import Cm



# Função para plotar o mapa
def plota_mapa(gridx, gridy, z, variavel, intervalos, cmap):
    plt.figure(figsize=(10, 6))
    grid_x, grid_y = np.meshgrid(gridx, gridy)
    norm = BoundaryNorm(intervalos, cmap.N)
    plt.contourf(grid_x, grid_y, z, levels=intervalos, cmap=cmap, norm=norm)
    plt.xlabel('Longitude')
    plt.ylabel('Latitude')
    plt.xlim(x_min, x_max)
    plt.ylim(y_min, y_max)
    plt.title(titulo)
    gdf.plot(ax=plt.gca(), edgecolor='black', facecolor='none', lw=1)
    cbar = plt.colorbar(label=legenda_eixo)
    cbar.ax.yaxis.set_ticks(intervalos)  # Define os intervalos explicitamente
    cbar.ax.yaxis.set_ticklabels(intervalos)  # Define as labels explicitamente

    plt.savefig(figura, dpi=300, bbox_inches="tight")
    plt.close()
    return figura

# Criar apresentação PPTX
def cria_pptx(imagem_path):
    prs = Presentation()
    prs.slide_width = Cm(29.7)
    prs.slide_height = Cm(21.0)
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    left = Cm(3)
    top = Cm(3)
    slide.shapes.add_picture(imagem_path, left, top, width=Cm(20), height=Cm(15))
    arquivo_pptx = "interpolacao_Kriging_exponential.pptx"
    prs.save(arquivo_pptx)
    print("Apresentação PowerPoint criada com sucesso!")

# Configurações iniciais
shapefile_path = '../00_STUFF/BAHIA.shp'
gdf = gpd.read_file(shapefile_path)

# Definição pelos limites do shapefile
x_min, y_min, x_max, y_max = gdf.total_bounds
res = 0.1
gridx, gridy = np.arange(x_min, x_max, res), np.arange(y_min, y_max, res)



# Intervalos e colormap
intervalos = [200, 400, 600, 800, 1000, 1200, 1400, 1600, 1800, 2000, 2200, 2400, 2600]
cmap = ListedColormap([
    'darkred', 'red', 'orange', 'yellow', 'greenyellow', 'limegreen',
    'green', 'cyan', 'deepskyblue', 'blue', 'mediumblue', 'darkblue'
])

# Executando o código
titulo="CLIMATOLOGIA - 1911 a 2024 (inventario) - Media anual"
figura="climatologia_inventario.png" 
legenda_eixo="Média mm/ano"
arquivo_input_nc = "./MEDIA/bahia_inventario__interpolacao_Kriging_exponential.nc"
# Carregar arquivo NetCDF
data = xr.open_dataset(arquivo_input_nc)
variavel = list(data.data_vars.keys())[0]  # Obtém a primeira variável disponível
z = data[variavel].values
imagem_gerada = plota_mapa(gridx, gridy, z, variavel, intervalos, cmap)



# Executando o código
titulo="CLIMATOLOGIA - 1981 a 2010 (Historico) - Média Anual"
figura="climatologia_1981a2010.png" 
legenda_eixo="Média mm/ano"
arquivo_input_nc = "./MEDIA/bahia_1981a2010__interpolacao_Kriging_exponential.nc"
# Carregar arquivo NetCDF
data = xr.open_dataset(arquivo_input_nc)
variavel = list(data.data_vars.keys())[0]  # Obtém a primeira variável disponível
z = data[variavel].values
imagem_gerada = plota_mapa(gridx, gridy, z, variavel, intervalos, cmap)


# Executando o código
titulo="CLIMATOLOGIA - 1991 a 2020 - Média Anual"
figura="climatologia_1991a2010.png" 
legenda_eixo="Média mm/ano"
arquivo_input_nc = "./MEDIA/bahia_1991a2020__interpolacao_Kriging_exponential.nc"
# Carregar arquivo NetCDF
data = xr.open_dataset(arquivo_input_nc)
variavel = list(data.data_vars.keys())[0]  # Obtém a primeira variável disponível
z = data[variavel].values
imagem_gerada = plota_mapa(gridx, gridy, z, variavel, intervalos, cmap)


cores = [
    (1.0, 0.0, 0.0),  # Vermelho
    (1.0, 0.3, 0.0),  # Laranja avermelhado
    (1.0, 0.5, 0.0),  # Laranja
    (1.0, 0.7, 0.0),  # Laranja amarelado
    (1.0, 1.0, 0.0),  # Amarelo
    (0.8, 1.0, 0.2),  # Verde amarelado
    (0.6, 1.0, 0.4),  # Verde claro
    (0.0, 1.0, 0.0),  # Verde
    (0.0, 0.8, 0.4),  # Verde escuro
    #(0.3, 1.0, 0.6),  # Verde água
    'seagreen',
    (0.0, 0.6, 0.8),  # Azul esverdeado
    (0.0, 0.4, 1.0),  # Azul claro
    (0.0, 0.2, 1.0),  # Azul médio
    (0.0, 0.0, 1.0),  # Azul
    (0.3, 0.0, 1.0),  # Azul arroxeado
    (0.5, 0.0, 1.0),  # Roxo azulado
    (0.7, 0.0, 1.0),  # Roxo médio
    (0.8, 0.2, 1.0),  # Roxo claro
    (0.9, 0.4, 1.0),  # Lilás
    (1.0, 0.6, 1.0)   # Rosa arroxeado
]

cmap = ListedColormap(cores)
norm = BoundaryNorm(intervalos, cmap.N)


# Definição dos intervalos e cores correspondentes
intervalos = [10, 12,14,16, 18, 20, 22, 24, 26, 28, 30, 32, 34, 36, 38, 40, 42, 44, 46, 48, 50]

# Executando o código
titulo="CLIMATOLOGIA - EVENTOS FORTES/SEVEROS - P90 médio"
figura="climatologia_P90_1991a2010.png" 
legenda_eixo="Média P90 (mm)"
arquivo_input_nc = "./P90/bahia_1991a2020__interpolacao_Kriging_exponential.nc"
# Carregar arquivo NetCDF
data = xr.open_dataset(arquivo_input_nc)
variavel = list(data.data_vars.keys())[0]  # Obtém a primeira variável disponível
z = data[variavel].values
imagem_gerada = plota_mapa(gridx, gridy, z, variavel, intervalos, cmap)

intervalos = [40, 45, 50, 55 , 60 , 65, 70, 72, 74, 76, 78, 80, 82, 84, 86, 88, 90, 94, 96, 98, 100]
# Executando o código
titulo="CLIMATOLOGIA - EVENTOS EXTREMOS - P99 médio"
figura="climatologia_P99_1991a2010.png" 
legenda_eixo="Média P99 (mm)"
arquivo_input_nc = "./P99/bahia_1991a2020__interpolacao_Kriging_exponential.nc"
# Carregar arquivo NetCDF
data = xr.open_dataset(arquivo_input_nc)
variavel = list(data.data_vars.keys())[0]  # Obtém a primeira variável disponível
z = data[variavel].values
imagem_gerada = plota_mapa(gridx, gridy, z, variavel, intervalos, cmap)


